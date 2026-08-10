import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0, 51, 102)
BLUE_LIGHT = RGBColor(173, 216, 230)
GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]
    
    title_text = title_placeholder.text_frame
    title_text.text = title
    for p in title_text.paragraphs:
        p.font.size = Pt(36)
        p.font.color.rgb = BLUE
        p.font.bold = True
    
    subtitle_text = subtitle_placeholder.text_frame
    subtitle_text.text = subtitle
    for p in subtitle_text.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY

def add_section_slide(title):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLUE
    
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

def add_text_slide(title, content_list):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_text = title_box.text_frame
    p = title_text.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    content_text = content_box.text_frame
    content_text.word_wrap = True
    
    for item in content_list:
        p = content_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.level = 0

def add_chart_slide(title, chart_path, notes=""):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_text = title_box.text_frame
    p = title_text.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.font.bold = True
    
    img_path = os.path.join("charts", chart_path)
    if os.path.exists(img_path):
        left = Inches(1)
        top = Inches(1.2)
        width = Inches(11)
        height = Inches(5.5)
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    else:
        content_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(2))
        content_text = content_box.text_frame
        p = content_text.add_paragraph()
        p.text = f"图表文件不存在: {chart_path}"
        p.font.size = Pt(18)
    
    if notes:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame
        notes_text.text = notes

def add_text_chart_slide(title, content_list, chart_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_text = title_box.text_frame
    p = title_text.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(5.5))
    content_text = content_box.text_frame
    content_text.word_wrap = True
    
    for item in content_list:
        p = content_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.level = 0
    
    img_path = os.path.join("charts", chart_path)
    if os.path.exists(img_path):
        left = Inches(6.2)
        top = Inches(1.2)
        width = Inches(6.6)
        height = Inches(5.5)
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    else:
        chart_box = slide.shapes.add_textbox(Inches(6.5), Inches(3), Inches(6), Inches(2))
        chart_text = chart_box.text_frame
        p = chart_text.add_paragraph()
        p.text = f"图表文件不存在: {chart_path}"
        p.font.size = Pt(14)

add_title_slide("AP Microeconomics", "Unit 1-Unit 3 核心内容精讲")

add_section_slide("Unit 1: Basic Economic Concepts")

add_text_slide("1.1 Scarcity & Choice 稀缺性与选择", [
    "Scarcity: 资源有限但欲望无限",
    "Three Fundamental Questions:",
    "  • What to produce? 生产什么?",
    "  • How to produce? 如何生产?",
    "  • For whom to produce? 为谁生产?",
    "Economic Thinking: 权衡取舍、边际分析、机会成本"
])

add_text_slide("1.2 Opportunity Cost 机会成本", [
    "Definition: 放弃的次优选项的价值",
    "公式: Opportunity Cost = 被放弃的次优选项的价值",
    "Key Point: 天下没有免费的午餐",
    "例: 读大学的机会成本 = 四年的工资收入"
])

add_chart_slide("1.3 Production Possibilities Curve (PPC)", "u1_ppc.png", 
                "A/B点: 有效率(Efficient)\nC点: 无效率(Inefficient/Unemployment)\nD点: 不可行(Unattainable)\n斜率 = 机会成本")

add_chart_slide("1.3.1 PPC & Economic Growth", "u1_ppc_growth.png",
                "PPC外移 = 经济增长\n原因: 技术进步、资源增加\nGrowth箭头指示增长方向")

add_text_slide("1.4 Marginal Utility 边际效用", [
    "Total Utility: 总效用",
    "Marginal Utility: 额外一单位商品带来的效用增量",
    "Law of Diminishing Marginal Utility:",
    "  随着消费数量增加，MU递减",
    "MU = ΔTU / ΔQ"
])

add_chart_slide("1.4.1 Diminishing Marginal Utility", "u1_dim_mu.png",
                "MU曲线向右下方倾斜\n每增加一单位，MU递减\nMU > 0: TU递增\nMU = 0: TU最大")

add_text_slide("1.5 Budget Line & Indifference Curve", [
    "Budget Line: Px*X + Py*Y = I",
    "斜率 = -Px/Py (机会成本)",
    "Indifference Curve (IC): 效用相同的商品组合",
    "MRS = -ΔY/ΔX = MUx/MUy",
    "Optimal Bundle: MRS = Px/Py"
])

add_chart_slide("1.5.1 Optimal Consumption Bundle", "u1_budget_ic.png",
                "切点E* = 最优消费束\nMRS = Px/Py\nIC与预算线相切\n在该点，消费者效用最大化")

add_section_slide("Unit 2: Supply and Demand")

add_text_slide("2.1 Demand 需求", [
    "Law of Demand: P↑ → Qd↓ (反向关系)",
    "需求曲线向右下方倾斜",
    "需求量变化(ΔQd): 沿曲线移动(价格变化)",
    "需求变化(ΔD): 曲线平移(非价格因素)",
    "非价格因素: 收入、偏好、相关商品价格、预期、人口"
])

add_chart_slide("2.1.1 Demand: Movement vs Shift", "u2_demand_move.png",
                "ΔP → 沿曲线移动(Movement)\nΔ非价格因素 → 曲线平移(Shift)\nD→D': 需求增加\nD→D'': 需求减少")

add_text_slide("2.2 Supply 供给", [
    "Law of Supply: P↑ → Qs↑ (正向关系)",
    "供给曲线向右上方倾斜",
    "供给量变化(ΔQs): 沿曲线移动",
    "供给变化(ΔS): 曲线平移",
    "非价格因素: 生产成本、技术、预期、天气、供给者数量"
])

add_chart_slide("2.2.1 Supply: Movement vs Shift", "u2_supply_move.png",
                "ΔP → 沿曲线移动\nΔ非价格因素 → 曲线平移\nS→S': 供给增加\nS→S'': 供给减少")

add_chart_slide("2.3 Market Equilibrium", "u2_equilibrium.png",
                "交点E = 均衡点\nPe = 均衡价格\nQe = 均衡数量\n短缺(Shortage): P < Pe\n过剩(Surplus): P > Pe")

add_text_chart_slide("2.3.1 Equilibrium Changes", 
                     ["需求增加 → Pe↑, Qe↑", "供给减少 → Pe↑, Qe↓"],
                     "u2_demand_increase.png")

add_chart_slide("2.3.2 Supply Decrease", "u2_supply_decrease.png",
                "供给减少 → 曲线左移\nPe↑, Qe↓\n新均衡点E1")

add_text_slide("2.5 Price Elasticity of Demand (PED)", [
    "PED = |%ΔQd / %ΔP|",
    "Midpoint Formula: PED = |ΔQ/Qavg / ΔP/Pavg|",
    "五种弹性类型:",
    "  • Perfectly Inelastic (PED=0)",
    "  • Inelastic (0<PED<1)",
    "  • Unit Elastic (PED=1)",
    "  • Elastic (1<PED<∞)",
    "  • Perfectly Elastic (PED=∞)"
])

add_chart_slide("2.5.1 Five Types of Elasticity", "u2_elasticity_five.png",
                "越陡峭越缺乏弹性\n越平坦越富有弹性\n垂直直线: PED=0\n水平直线: PED=∞")

add_text_slide("2.6 Price Controls 价格管制", [
    "Price Ceiling (价格上限): 最高限价",
    "  • 必须低于Pe才有效",
    "  • 导致短缺(Shortage)",
    "  • 例: 租金管制",
    "Price Floor (价格下限): 最低限价",
    "  • 必须高于Pe才有效",
    "  • 导致过剩(Surplus)",
    "  • 例: 最低工资"
])

add_chart_slide("2.6.1 Price Ceiling", "u2_price_ceiling.png",
                "Pc < Pe\n短缺量 = Qd - Qs\n非价格配给机制出现")

add_chart_slide("2.6.2 Price Floor", "u2_price_floor.png",
                "Pf > Pe\n过剩量 = Qs - Qd")

add_chart_slide("2.7 Tax Incidence 税收归宿", "u2_tax_incidence.png",
                "税收使供给曲线上移\n新均衡点E1\n消费者支付Pc > Pe\n生产者获得Pp < Pe\n税负 = Pc - Pp\n归宿取决于弹性")

add_chart_slide("2.7.1 Per-Unit Subsidy", "u2_subsidy.png",
                "补贴使供给曲线下移\n新均衡点E1\n消费者支付Pc < Pe\n生产者获得Pp > Pe\n补贴额 = Pp - Pc\nDWL = 无谓损失")

add_chart_slide("2.8 Consumer & Producer Surplus", "u2_cs_ps.png",
                "CS = 消费者剩余(需求曲线下方)\nPS = 生产者剩余(供给曲线上方)\n总剩余 = CS + PS\nE点 = 社会最优")

add_text_slide("2.9 International Trade 国际贸易", [
    "Free Trade: 自由贸易",
    "  • 进口国: P = Pw < Pe",
    "  • 出口国: P = Pw > Pe",
    "Tariff: 关税(进口税)",
    "  • 提高国内价格",
    "  • 减少进口量",
    "  • 产生DWL",
    "Quota: 配额(进口限制)",
    "  • 限制进口数量",
    "  • 产生Quota Rent"
])

add_chart_slide("2.9.1 Tariff", "u2_tariff.png",
                "Pw+t > Pw\n进口减少\nGov Revenue = 关税收入\nDWL = 生产扭曲 + 消费扭曲")

add_chart_slide("2.9.2 Quota", "u2_quota.png",
                "进口量限制为Quota\n价格上升到Pq\nQuota Rent = 配额租金\nDWL = 与关税类似")

add_section_slide("Unit 3: Production and Costs")

add_text_slide("3.1 Production Function", [
    "TP (Total Product): 总产量",
    "AP (Average Product): AP = TP/L",
    "MP (Marginal Product): MP = ΔTP/ΔL",
    "Law of Diminishing Marginal Returns:",
    "  随着可变投入增加，MP最终递减",
    "MP > AP: AP递增",
    "MP = AP: AP最大",
    "MP < AP: AP递减"
])

add_chart_slide("3.1.1 TP, MP & AP", "u3_tp_mp_ap.png",
                "MP曲线先升后降\nMP穿过AP最高点\nTP斜率 = MP")

add_text_slide("3.2 Costs of Production", [
    "Fixed Costs (TFC): 固定成本(不随产量变化)",
    "Variable Costs (TVC): 可变成本(随产量变化)",
    "Total Costs: TC = TFC + TVC",
    "Marginal Cost: MC = ΔTC/ΔQ",
    "Average Costs:",
    "  • AFC = TFC/Q (持续下降)",
    "  • AVC = TVC/Q (U形)",
    "  • ATC = TC/Q = AFC + AVC (U形)"
])

add_chart_slide("3.2.1 Total Cost Curves", "u3_total_cost.png",
                "TFC = 水平直线\nTVC = 从原点出发的曲线\nTC = TFC + TVC\nTC与TVC平行")

add_chart_slide("3.3 Unit Cost Curves", "u3_unit_cost.png",
                "ATC & AVC = U形\nAFC = 双曲线(持续下降)\nMC穿过AVC和ATC最低点\nATC = AVC + AFC")

add_text_slide("3.4 Long-Run Average Total Cost", [
    "LRATC = 长期平均总成本",
    "由SRATC曲线的包络线构成",
    "Economies of Scale: 规模经济(LRATC↓)",
    "Constant Returns to Scale: 规模报酬不变",
    "Diseconomies of Scale: 规模不经济(LRATC↑)",
    "Minimum Efficient Scale (MES): LRATC最低点"
])

add_chart_slide("3.4.1 LRATC Curve", "u3_lratc.png",
                "LRATC = SRATC包络线\n三种规模状态:\n1. Economies of Scale\n2. Constant Returns\n3. Diseconomies of Scale")

add_text_slide("3.5 Perfect Competition 完全竞争市场", [
    "四个特征:",
    "  • 大量买者卖者",
    "  • 产品同质",
    "  • 自由进出",
    "  • 价格接受者",
    "Demand Curve: D = MR = AR = P (水平直线)",
    "Profit Max: MR = MC",
    "短期均衡三种情况:",
    "  • 经济利润 > 0",
    "  • 经济利润 = 0",
    "  • 经济利润 < 0 (亏损)"
])

add_chart_slide("3.5.1 Short-Run Economic Profit", "u3_pc_profit.png",
                "P > ATC\nProfit = (P - ATC) * Q*\n绿色阴影 = 利润面积")

add_chart_slide("3.5.2 Short-Run Loss", "u3_pc_loss.png",
                "P < ATC\nLoss = (ATC - P) * Q*\n红色阴影 = 亏损面积")

add_text_slide("3.5.3 Shutdown Decision 停产决策", [
    "继续经营: P ≥ AVC (可覆盖部分FC)",
    "停产: P < AVC (无法覆盖VC)",
    "Shutdown Point: P = min(AVC)",
    "短期: FC沉没成本，不予考虑",
    "长期: 所有成本可变，退出市场"
])

add_chart_slide("3.5.4 Shutdown Point", "u3_pc_shutdown.png",
                "P = min(AVC)\n在该点，继续经营和停产损失相同\n均等于TFC")

add_text_slide("3.5.5 Long-Run Equilibrium", [
    "长期均衡条件: P = MR = MC = ATC",
    "经济利润 = 0",
    "厂商数量调整:",
    "  • 正利润 → 新厂商进入 → S↑ → P↓",
    "  • 负利润 → 厂商退出 → S↓ → P↑",
    "最终达到零经济利润"
])

add_chart_slide("3.5.5 Long-Run Equilibrium", "u3_pc_lr_equilibrium.png",
                "P = MR = MC = min(ATC)\n零经济利润\nAllocatively Efficient (P=MC)\nProductively Efficient (P=minATC)")

output_path = "AP微观经济学讲义_Unit1-3.pptx"
prs.save(output_path)
print(f"PPT文件已生成: {output_path}")